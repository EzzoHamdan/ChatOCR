from flask import Flask, request, jsonify, render_template
from werkzeug.utils import secure_filename
from openai import OpenAI
from docx import Document
from pptx import Presentation
import os
from dotenv import load_dotenv
import pytesseract
from PIL import Image

load_dotenv()
client = OpenAI(
    api_key=os.environ.get("OPENAI_API_KEY"),
)

# Flask app
app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.secret_key = 'your_secret_key'
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # Limit 16 MB

# Function to extract text from image
def extract_image(image_path):
    image = Image.open(image_path)
    return pytesseract.image_to_string(image)

# Function to extract text from pdf
def extract_pdf(pdf_path, pages_input):
    from PyPDF2 import PdfReader
    
    def parse_page_range(pages_input):
        pages = set()
        if not pages_input:
            return None  # Process all pages if no input
        for part in pages_input.split(','):
            if '-' in part:
                start, end = map(int, part.split('-'))
                pages.update(range(start - 1, end))  # Page numbers are 0-based
            else:
                pages.add(int(part) - 1)  # Single page
        return sorted(pages)

    pages_to_extract = parse_page_range(pages_input)
    text = ""
    with open(pdf_path, 'rb') as f:
        reader = PdfReader(f)
        total_pages = len(reader.pages)
        
        if pages_to_extract:
            pages_to_extract = [p for p in pages_to_extract if p < total_pages]  # Avoid invalid pages
        else:
            pages_to_extract = range(total_pages)  # Default: All pages
        
        for i in pages_to_extract:
            text += reader.pages[i].extract_text() or ""  # Extract text safely
    return text


# Function to extract text from doc
def extract_docx(docx_path, pages_input):
    doc = Document(docx_path)
    paragraphs = doc.paragraphs
    
    if pages_input:
        start, end = map(int, pages_input.split('-'))
        paragraphs = paragraphs[start-1:end]  # Slice paragraphs based on range
    
    text = "\n".join(para.text for para in paragraphs)
    return text

# Function to extract text from ppt
def extract_pptx(pptx_path, pages_input):
    def parse_page_range(pages_input):
        pages = set()
        if not pages_input:
            return None
        for part in pages_input.split(','):
            if '-' in part:
                start, end = map(int, part.split('-'))
                pages.update(range(start - 1, end))
            else:
                pages.add(int(part) - 1)
        return sorted(pages)

    pages_to_extract = parse_page_range(pages_input)
    text = ""
    presentation = Presentation(pptx_path)
    total_slides = len(presentation.slides)
    
    if pages_to_extract:
        pages_to_extract = [p for p in pages_to_extract if p < total_slides]
    else:
        pages_to_extract = range(total_slides)
    
    for i in pages_to_extract:
        slide = presentation.slides[i]
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + "\n"
    return text

# Function to generate response from ChatGPT
def chatgpt_response(prompt):
    response = client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="gpt-3.5-turbo"
    )
    return response.choices[0].message.content

# Route for the homepage
@app.route('/')
def home():
    return render_template('index.html')

# Route to handle file uploads and queries
@app.route('/process', methods=['POST'])
def process():
    try:
        query = request.form.get('query')
        file = request.files.get('file')
        pages_input = request.form.get('pages')  # Fetch user-specified pages
        
        if not query and not file:
            return jsonify({'error': 'Please provide a query or file.'}), 400

        combined_text = query or ""
        task = request.form.get('task')
        prompt_to_send = ""
        extracted_text = ""  # to prevent uninitialized use

        # Parse task
        if task == 'summarize':
            prompt_to_send = f"Summarize this text:\n{combined_text}"
        elif task == 'question':
            prompt_to_send = f"Answer the following question based on the text:\n{combined_text}\n\nQuestion: {query}"
        elif task == 'rewrite':
            prompt_to_send = f"Rephrase the following text:\n{combined_text}"
        elif task == 'explain':
            prompt_to_send = f"Explain the following text for an average person:\n{combined_text}"
        elif task == 'explainplus':
            prompt_to_send = f"Explain the following text for a dummy:\n{combined_text}"
        elif task == 'inquiry':
            prompt_to_send = f"Follow the instructions given with the text:\n{combined_text}"

        if file:
            filename = secure_filename(file.filename)
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(file_path)

            # Extract text based on file type
            if filename.endswith('.pdf'):
                extracted_text = extract_pdf(file_path, pages_input)
            elif filename.endswith('.docx'):
                extracted_text = extract_docx(file_path, pages_input)
            elif filename.endswith('.pptx'):
                extracted_text = extract_pptx(file_path, pages_input)
            elif filename.endswith(('.png', '.jpg', '.jpeg')):
                extracted_text = extract_image(file_path)
            else:
                return jsonify({'error': 'Unsupported file type.'}), 400

            combined_text += "\n" + extracted_text
            os.remove(file_path)  # Clean up uploaded file

        # Get response from ChatGPT
        reply = chatgpt_response(prompt_to_send + "\n" + combined_text)
        return jsonify({'response': reply, 'prompt': prompt_to_send + "\n" + extracted_text})

    except Exception as e:
        return jsonify({'error': str(e)}), 500

# Run the Flask app
if __name__ == '__main__':
    app.run(debug=True)
