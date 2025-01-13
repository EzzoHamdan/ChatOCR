# utils/file_extraction.py

import pytesseract
from PIL import Image
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

def parse_page_range(pages_input):
    pages = set()
    if not pages_input:
        return None  # Process all pages if no input
    for part in pages_input.split(','):
        if '-' in part:
            start, end = map(int, part.split('-'))
            pages.update(range(start - 1, end))  # Page numbers are 0-based
        else:
            pages.add(int(part) - 1)  # Single page
    return sorted(pages)

def extract_image(image_path):
    image = Image.open(image_path)
    text = pytesseract.image_to_string(image)
    return text

# Function to extract text from pdf
def extract_pdf(pdf_path, pages_input):
    pages_to_extract = parse_page_range(pages_input)
    text = ""
    with open(pdf_path, 'rb') as f:
        reader = PdfReader(f)
        total_pages = len(reader.pages)
        
        if pages_to_extract:
            pages_to_extract = [p for p in pages_to_extract if p < total_pages]  # Avoid invalid pages
        else:
            pages_to_extract = range(total_pages)  # Default: All pages
        
        for i in pages_to_extract:
            text += reader.pages[i].extract_text() or ""  # Extract text safely
    return text

def extract_docx(docx_path, pages_input):
    doc = Document(docx_path)
    paragraphs = doc.paragraphs
    
    if pages_input:
        start, end = map(int, pages_input.split('-'))
        paragraphs = paragraphs[start-1:end]  # Slice paragraphs based on range
    
    text = "\n".join(para.text for para in paragraphs)
    return text

# Function to extract text from ppt
def extract_pptx(pptx_path, pages_input):
    pages_to_extract = parse_page_range(pages_input)
    text = ""
    presentation = Presentation(pptx_path)
    total_slides = len(presentation.slides)
    
    if pages_to_extract:
        pages_to_extract = [p for p in pages_to_extract if p < total_slides]
    else:
        pages_to_extract = range(total_slides)
    
    for i in pages_to_extract:
        slide = presentation.slides[i]
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + "\n"
    return text
